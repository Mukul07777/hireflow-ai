"""Builds the NexFlow AI presentation.  Run:  python3 scripts/build_deck.py"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SHOTS = os.path.join(ROOT, "docs", "screenshots")

# palette taken from the product itself
NAVY, TEAL, GOLD = RGBColor(0x0D, 0x3B, 0x4F), RGBColor(0x1C, 0x7A, 0x93), RGBColor(0xB8, 0x89, 0x4A)
GOLD_LT = RGBColor(0xE0, 0xB4, 0x63)
CREAM, INK, MUTED = RGBColor(0xFB, 0xF6, 0xEC), RGBColor(0x21, 0x1E, 0x19), RGBColor(0x6B, 0x63, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE4, 0xDA, 0xC6)
ICE = RGBColor(0xC9, 0xD6, 0xDC)
STEEL = RGBColor(0x9F, 0xB4, 0xBD)
ROSE_BG, ROSE_LN, ROSE_TX = RGBColor(0xFB, 0xEA, 0xF0), RGBColor(0xED, 0x93, 0xB1), RGBColor(0x99, 0x35, 0x56)
AMB_BG, AMB_TX = RGBColor(0xFA, 0xEE, 0xDA), RGBColor(0x85, 0x4F, 0x0B)
MINT_BG, MINT_LN, MINT_TX = RGBColor(0xE1, 0xF5, 0xEE), RGBColor(0x9F, 0xE1, 0xCB), RGBColor(0x0F, 0x6E, 0x56)
PANEL = RGBColor(0x15, 0x4A, 0x5E)
SAND = RGBColor(0xEF, 0xE7, 0xD6)

H, B = "Cambria", "Calibri"

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(bg=CREAM):
    s = prs.slides.add_slide(BLANK)
    bgfill = s.background.fill
    bgfill.solid()
    bgfill.fore_color.rgb = bg
    return s


def noline(shp):
    shp.line.fill.background()


def rrect(s, x, y, w, h, fill, line=None, radius=0.08):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        noline(shp)
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    shp.shadow.inherit = False
    if shp.has_text_frame:
        shp.text_frame.text = ""
    return shp


def oval(s, x, y, w, h, fill, line=None):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        noline(shp)
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1.5)
    shp.shadow.inherit = False
    return shp


def txt(s, text, x, y, w, h, size=14, bold=False, color=INK, font=B, align=PP_ALIGN.LEFT,
        italic=False, spacing=None, anchor=MSO_ANCHOR.TOP, shrink=False):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.line_spacing = spacing
        r = p.add_run(); r.text = line
        f = r.font
        f.name = font; f.size = Pt(size); f.bold = bold; f.italic = italic; f.color.rgb = color
    return tb


def bullets(s, items, x, y, w, h, size=13, color=INK):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(9); p.line_spacing = 1.25
        r = p.add_run(); r.text = "•  " + item
        r.font.name = B; r.font.size = Pt(size); r.font.color.rgb = color
    return tb


def pic(s, name, x, y, w, h):
    return s.shapes.add_picture(os.path.join(SHOTS, name), Inches(x), Inches(y), Inches(w), Inches(h))


def eyebrow(s, text, x=0.7, y=0.45, color=GOLD):
    txt(s, text, x, y, 9, 0.28, size=11.5, bold=True, color=color)


def title(s, text, y=0.8, size=34, w=11.9, color=INK, x=0.7, h=1.0):
    txt(s, text, x, y, w, h, size=size, bold=True, color=color, font=H, spacing=1.15)


def statcard(s, x, y, w, h, value, label, vcolor=NAVY, vsize=32, fill=WHITE, line=BORDER, lcolor=MUTED):
    rrect(s, x, y, w, h, fill, line, 0.12)
    txt(s, value, x, y + 0.14, w, 0.62, size=vsize, bold=True, color=vcolor, font=H, align=PP_ALIGN.CENTER)
    txt(s, label, x + 0.12, y + h - 0.62, w - 0.24, 0.55, size=10.5, color=lcolor, align=PP_ALIGN.CENTER, spacing=1.1)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def dark_decor(s):
    o = oval(s, 9.7, -1.9, 6.2, 6.2, RGBColor(0x14, 0x4A, 0x5F)); noline(o)
    o2 = oval(s, -1.6, 4.4, 4.4, 4.4, RGBColor(0x13, 0x44, 0x57)); noline(o2)


# ══ 1 · TITLE ═══════════════════════════════════════════════════════════════
s = slide(NAVY); dark_decor(s)
rrect(s, 0.9, 1.0, 1.0, 1.0, TEAL, None, 0.2)
txt(s, "NX", 0.9, 1.28, 1.0, 0.5, size=26, bold=True, color=WHITE, font=H, align=PP_ALIGN.CENTER)
txt(s, "NexFlow AI", 2.15, 1.0, 9, 1.05, size=54, bold=True, color=WHITE, font=H, anchor=MSO_ANCHOR.MIDDLE)
txt(s, "Six AI agents that share one memory.", 0.9, 2.4, 11, 0.6, size=26, italic=True, color=GOLD_LT, font=H)
txt(s, "A multi-agent business-intelligence platform for Indian SMBs — where hiring, sales,\nsupport and customer care finally recognise the same customer.",
    0.9, 3.15, 10.6, 1.0, size=14.5, color=ICE, spacing=1.45)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(4.35), Inches(3.2), Pt(2.5))
ln.fill.solid(); ln.fill.fore_color.rgb = GOLD; noline(ln); ln.shadow.inherit = False
txt(s, "Mukul  ·  IIIT-Delhi", 0.9, 4.62, 8, 0.35, size=16, bold=True, color=WHITE)
txt(s, "with Supriya Bansal  ·  IIT-Delhi", 0.9, 5.02, 8, 0.32, size=13, color=STEEL)
txt(s, "nexflow-ai-india.netlify.app", 0.9, 5.62, 8, 0.32, size=12.5, color=GOLD_LT)
notes(s, "Hook: most AI tools for SMBs are single-purpose, and none of them know the same customer.")

# ══ 2 · THE PROBLEM ═════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "THE PROBLEM")
title(s, "Six agents that don't share memory\nare still six silos.", y=0.82, size=36, h=1.6)
txt(s, "MONDAY", 0.7, 2.72, 5.3, 0.28, size=11, bold=True, color=GOLD)
rrect(s, 0.7, 3.05, 5.3, 1.3, WHITE, BORDER, 0.12)
txt(s, "A customer files an angry billing complaint.\nSupport handles it. Support knows.",
    0.95, 3.28, 4.85, 0.9, size=13.5, color=INK, spacing=1.35)
txt(s, "TUESDAY", 0.7, 4.6, 5.3, 0.28, size=11, bold=True, color=ROSE_TX)
rrect(s, 0.7, 4.93, 5.3, 1.3, ROSE_BG, ROSE_LN, 0.12)
txt(s, "Sales emails that same company an upsell\npitch. Sales had no idea.",
    0.95, 5.16, 4.85, 0.9, size=13.5, color=ROSE_TX, spacing=1.35)
rrect(s, 6.6, 2.72, 6.05, 3.5, NAVY, None, 0.12)
txt(s, "Nobody did anything wrong.", 6.95, 3.1, 5.4, 0.5, size=21, bold=True, color=WHITE, font=H)
txt(s, "The two teams simply cannot see each other.\n\nThat isn't a small annoyance — that's how you\nlose the account.",
    6.95, 3.72, 5.4, 1.7, size=14, color=ICE, spacing=1.4)
txt(s, "So we built a seventh system to fix exactly this.", 6.95, 5.6, 5.4, 0.4, size=13, bold=True, italic=True, color=GOLD_LT)
notes(s, "Make it concrete. This scenario is the reason the whole product exists.")

# ══ 3 · SEVEN SYSTEMS ═══════════════════════════════════════════════════════
s = slide()
eyebrow(s, "WHAT WE BUILT")
title(s, "Seven AI systems. One platform.", y=0.78, size=34, h=0.65)
txt(s, "6 agents  +  the Company Brain  =  7 AI", 0.7, 1.5, 6, 0.35, size=15, bold=True, color=TEAL)
pic(s, "home.png", 0.7, 1.95, 8.15, 4.05)
rows = [("SMB Brain", "WhatsApp chats to a real CRM"),
        ("HireFlow", "7-agent hiring pipeline"),
        ("SalesFlow", "Prospecting and outreach"),
        ("SupportFlow", "Knowledge base from your docs"),
        ("CareFlow", "Triage with human approval"),
        ("War Room", "All agents fire in parallel"),
        ("Company Brain", "Shared memory — the glue")]
for i, (name, desc) in enumerate(rows):
    y = 1.95 + i * 0.6
    last = i == len(rows) - 1
    oval(s, 9.0, y + 0.08, 0.34, 0.34, TEAL if last else SAND)
    txt(s, str(i + 1) if not last else "7", 9.0, y + 0.14, 0.34, 0.25, size=10.5, bold=True,
        color=WHITE if last else MUTED, align=PP_ALIGN.CENTER)
    txt(s, name, 9.48, y + 0.04, 2.0, 0.28, size=12.5, bold=True, color=TEAL if last else INK)
    txt(s, desc, 11.45, y + 0.06, 1.6, 0.4, size=9.5, color=MUTED, spacing=1.15)
txt(s, "Team is an admin feature, not an AI system — kept separate on purpose. We don't inflate the count.",
    0.7, 6.35, 11.95, 0.4, size=12.5, italic=True, color=MUTED)
notes(s, "Name them quickly. Team is admin, not AI — we don't inflate the count.")

# ══ 4 · THE CLOSED LOOP ═════════════════════════════════════════════════════
s = slide()
eyebrow(s, "THE IDEA")
title(s, "Every agent reads the shared memory\nbefore it acts.", y=0.78, size=32, h=1.3)
agents = ["HireFlow", "SalesFlow", "SupportFlow", "CareFlow", "SMB Brain", "War Room"]
for i, nm in enumerate(agents):
    x = 0.75 + (i % 3) * 1.6
    y = 2.5 + (i // 3) * 1.15
    rrect(s, x, y, 1.45, 0.92, WHITE, BORDER, 0.12)
    txt(s, nm, x, y + 0.32, 1.45, 0.3, size=10.5, bold=True, color=INK, align=PP_ALIGN.CENTER)
arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.5), Inches(3.42), Inches(1.0), Inches(0.3))
arrow.fill.solid(); arrow.fill.fore_color.rgb = GOLD; noline(arrow); arrow.shadow.inherit = False
oval(s, 6.65, 2.6, 2.2, 2.2, NAVY, TEAL)
txt(s, "COMPANY\nBRAIN", 6.65, 3.28, 2.2, 0.85, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, spacing=1.2)
facts = [("Identity resolution", "email to phone to name to account"),
         ("Unified timeline", "every action on one company-wide feed"),
         ("Next-best-action", "deterministic, ranked, explainable")]
for i, (t, d) in enumerate(facts):
    y = 2.55 + i * 0.8
    oval(s, 9.2, y + 0.04, 0.32, 0.32, GOLD)
    txt(s, str(i + 1), 9.2, y + 0.1, 0.32, 0.24, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, t, 9.68, y, 3.4, 0.3, size=13, bold=True, color=INK)
    txt(s, d, 9.68, y + 0.3, 3.4, 0.34, size=10.5, color=MUTED)
txt(s, "They write to it — and read from it before they draft.", 0.75, 5.35, 8, 0.35,
    size=13, bold=True, italic=True, color=TEAL)
notes(s, "The architecture in one picture. Emphasise: read BEFORE acting.")

# ══ 5 · DEMO ROADMAP ════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "WHAT YOU'RE ABOUT TO SEE")
title(s, "Live demo — four things", y=0.78, size=34, h=0.65)
steps = [("1", "CareFlow", "A real customer complaint lands. Care logs it.", TEAL),
         ("2", "SalesFlow", "Sales targets that same customer — and the Brain stops them.", GOLD),
         ("3", "Company Brain", "Two records become one. Ranked cross-team actions.", NAVY),
         ("4", "Proof", "The deterministic engine, run live in the terminal.", ROSE_TX)]
for i, (n, t, d, c) in enumerate(steps):
    y = 1.7 + i * 1.15
    rrect(s, 0.7, y, 11.95, 1.0, WHITE, BORDER, 0.1)
    oval(s, 0.95, y + 0.25, 0.5, 0.5, c)
    txt(s, n, 0.95, y + 0.35, 0.5, 0.32, size=15, bold=True, color=WHITE, font=H, align=PP_ALIGN.CENTER)
    txt(s, t, 1.68, y + 0.32, 2.7, 0.35, size=15, bold=True, color=c)
    txt(s, d, 4.4, y + 0.34, 8.0, 0.4, size=13, color=INK)
notes(s, "Tell them the shape of the demo before switching. Judges follow better when they know what's coming.")

# ══ 6 · DEMO DIVIDER ════════════════════════════════════════════════════════
s = slide(NAVY); dark_decor(s)
txt(s, "LIVE DEMO", 0.9, 2.25, 11.5, 1.3, size=64, bold=True, color=WHITE, font=H)
txt(s, "nexflow-ai-india.netlify.app", 0.9, 3.7, 9, 0.5, size=19, color=GOLD_LT)
rrect(s, 0.9, 4.6, 5.3, 0.62, TEAL, None, 0.3)
txt(s, "Switch to the browser now", 0.9, 4.77, 5.3, 0.35, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
notes(s, "SWITCH TO BROWSER. Order: CareFlow ticket, SalesFlow Outreach warnings, Company Brain merge, terminal scenario script. Then return to the next slide.")

# ══ 7 · RECAP ═══════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "WHAT JUST HAPPENED")
title(s, "Sales was stopped before it sent.", y=0.78, size=34, h=0.65)
pic(s, "closed-loop.png", 0.7, 1.7, 5.5, 4.7)
cards = [(1.7, ROSE_BG, ROSE_LN, ROSE_TX, "Exact contact match",
          "Priya is already known to Care — and has an\nopen complaint. Align before pitching."),
         (3.35, AMB_BG, GOLD, AMB_TX, "Account-level match",
          "Ananya is brand new — but her colleague has\nan open ticket at the same company."),
         (5.0, MINT_BG, MINT_LN, MINT_TX, "Net-new — and it says so",
          "When there is genuinely no history, it tells you.\nNo invented intelligence.")]
for y, bg, ln_, tx, head, body in cards:
    rrect(s, 6.6, y, 6.05, 1.45, bg, ln_, 0.12)
    txt(s, head, 6.88, y + 0.16, 5.5, 0.3, size=13.5, bold=True, color=tx)
    txt(s, body, 6.88, y + 0.52, 5.5, 0.8, size=12, color=tx, spacing=1.3)
notes(s, "The warning isn't just a badge — it's injected into the prompt, so the model writes a different email.")

# ══ 8 · COMPANY BRAIN ═══════════════════════════════════════════════════════
s = slide()
eyebrow(s, "THE SHARED MEMORY")
title(s, "Two records became one.", y=0.78, size=34, h=0.65)
pic(s, "company-brain.png", 0.7, 1.7, 7.6, 4.48)
statcard(s, 8.6, 1.7, 4.05, 1.3, "1", "duplicate identity merged\nacross Care and Sales", TEAL, 34)
statcard(s, 8.6, 3.15, 4.05, 1.3, "14", "people and businesses\nknown to the company", NAVY, 34)
statcard(s, 8.6, 4.6, 4.05, 1.58, "HIGH", "“Resolve the open complaint\nbefore sales pitches an upsell”",
         ROSE_TX, 24, ROSE_BG, ROSE_LN, ROSE_TX)
notes(s, "No single agent could produce that recommendation. It exists only because the memory is shared.")

# ══ 9 · DETERMINISTIC BY DESIGN ═════════════════════════════════════════════
s = slide()
eyebrow(s, "ENGINEERING DECISION")
title(s, "The Brain decides. The model drafts.", y=0.78, size=34, h=0.65)
rrect(s, 0.7, 1.75, 5.95, 3.95, NAVY, None, 0.12)
txt(s, "Deterministic core", 1.0, 2.0, 5.3, 0.4, size=16, bold=True, color=GOLD_LT)
bullets(s, ["Identity resolution and every recommendation are pure functions",
            "Same input produces identical output, every run",
            "Fuzzy matching: “Raj Patel” = “Rajesh Patel” at one domain — two different colleagues never merge",
            "Reproducible in 2 seconds, no API key needed"],
        1.0, 2.5, 5.3, 3.0, 12.5, ICE)
rrect(s, 6.95, 1.75, 5.7, 3.95, WHITE, BORDER, 0.12)
txt(s, "Probabilistic layer", 7.25, 2.0, 5.1, 0.4, size=16, bold=True, color=TEAL)
bullets(s, ["Llama 3.3 70B via Groq writes the actual copy",
            "It receives the Brain's facts as injected context",
            "A human approves before anything is sent",
            "If the model is down, the safety check still fires"],
        7.25, 2.5, 5.1, 3.0, 12.5, INK)
txt(s, "That separation is why a rate-limited model cannot cause a bad business decision.",
    0.7, 5.9, 11.95, 0.4, size=13.5, bold=True, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
notes(s, "This answers 'is it just a ChatGPT wrapper'. No — the decisions are deterministic.")

# ══ 10 · MEASURED RUN (CHART) ═══════════════════════════════════════════════
s = slide()
eyebrow(s, "PROOF, NOT CLAIMS")
title(s, "Measured against the live model.", y=0.78, size=34, h=0.65)
cd = CategoryChartData()
cd.categories = ["JD analysis", "Bias audit", "Sales email\n+ Brain", "Care reply\n+ Brain"]
cd.add_series("Latency (ms)", (433, 1142, 1920, 916))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), Inches(1.65),
                        Inches(7.6), Inches(4.55), cd)
ch = gf.chart
ch.has_legend = False
ch.has_title = True
ch.chart_title.text_frame.text = "Real Groq API run — latency per step"
ctf = ch.chart_title.text_frame.paragraphs[0].runs[0].font
ctf.size, ctf.name, ctf.color.rgb, ctf.bold = Pt(12), B, MUTED, False
plot = ch.plots[0]
plot.has_data_labels = True
dl = plot.data_labels
dl.number_format, dl.number_format_is_linked = '0', False
dl.position = XL_LABEL_POSITION.OUTSIDE_END
dl.font.size, dl.font.name, dl.font.color.rgb = Pt(10), B, INK
for pt_i, colr in enumerate([TEAL, TEAL, GOLD, TEAL]):
    p = plot.series[0].points[pt_i]
    p.format.fill.solid(); p.format.fill.fore_color.rgb = colr
for ax in (ch.category_axis, ch.value_axis):
    ax.tick_labels.font.size = Pt(10); ax.tick_labels.font.name = B
    ax.tick_labels.font.color.rgb = MUTED
ch.value_axis.has_major_gridlines = True
statcard(s, 8.6, 1.65, 4.05, 1.42, "4.4 s", "full four-call run", NAVY, 32)
statcard(s, 8.6, 3.22, 4.05, 1.42, "₹0.096", "measured cost of this run", TEAL, 32)
statcard(s, 8.6, 4.79, 4.05, 1.42, "128", "tests · 11 suites · CI on every push", GOLD, 32)
txt(s, "Reproducible:  node scripts/real-run.mjs  — writes the log straight into the README.",
    0.7, 6.45, 11.95, 0.35, size=12, italic=True, color=MUTED)
notes(s, "Every number was measured, not estimated. The script is in the repo.")

# ══ 11 · COST (CHART) ═══════════════════════════════════════════════════════
s = slide()
eyebrow(s, "THE QUESTION EVERY SMB ASKS")
title(s, "What does it actually cost to run?", y=0.78, size=34, h=0.65)
cd2 = CategoryChartData()
cd2.categories = ["100 runs", "500 runs", "1,000 runs", "5,000 runs"]
cd2.add_series("Monthly compute (Rs)", (19, 95, 190, 950))
gf2 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), Inches(1.7),
                         Inches(7.4), Inches(4.5), cd2)
c2 = gf2.chart
c2.has_legend = False
c2.has_title = True
c2.chart_title.text_frame.text = "Monthly Groq compute by usage volume (Rs)"
t2 = c2.chart_title.text_frame.paragraphs[0].runs[0].font
t2.size, t2.name, t2.color.rgb, t2.bold = Pt(12), B, MUTED, False
pl2 = c2.plots[0]
pl2.has_data_labels = True
d2 = pl2.data_labels
d2.number_format, d2.number_format_is_linked = '0', False
d2.position = XL_LABEL_POSITION.OUTSIDE_END
d2.font.size, d2.font.name, d2.font.color.rgb = Pt(10), B, INK
for pt_i, colr in enumerate([TEAL, TEAL, GOLD, NAVY]):
    p = pl2.series[0].points[pt_i]
    p.format.fill.solid(); p.format.fill.fore_color.rgb = colr
for ax in (c2.category_axis, c2.value_axis):
    ax.tick_labels.font.size = Pt(10); ax.tick_labels.font.name = B
    ax.tick_labels.font.color.rgb = MUTED
rrect(s, 8.45, 1.7, 4.2, 4.5, NAVY, None, 0.12)
txt(s, "₹0.19", 8.45, 2.45, 4.2, 0.85, size=44, bold=True, color=GOLD_LT, font=H, align=PP_ALIGN.CENTER)
txt(s, "per full HireFlow pipeline run\n(8 Groq calls)", 8.7, 3.05, 3.7, 0.7,
    size=12, color=ICE, align=PP_ALIGN.CENTER, spacing=1.3)
txt(s, "Computed from Groq's published pricing\nand our own token tracking — not an\ninvented marketing figure.",
    8.7, 4.6, 3.7, 1.1, size=11, color=STEEL, align=PP_ALIGN.CENTER, spacing=1.3)
notes(s, "Rs 190 a month at genuinely high usage. That's the number an SMB owner actually asks about.")

# ══ 12 · INDIA-FIRST ════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "NOT A WESTERN TOOL WITH A RUPEE SIGN")
title(s, "Built for how India actually works.", y=0.78, size=34, h=0.65)
india = [("Compliance by Indian law", "Deterministic JD checks written against the Equal Remuneration Act, POSH Act and Code on Wages."),
         ("PAN · GSTIN · Udyam", "Structural and checksum validation, verified against the standard public GSTIN vector."),
         ("Hinglish, not translation", "The AI answers in natural Roman-script Hinglish, the way Indian professionals write."),
         ("WhatsApp-native", "Tier-2 businesses live on WhatsApp. SMB Brain turns raw chats into a structured CRM."),
         ("DPDP Act, 2023", "A real data-subject deletion function, with cascade behaviour documented."),
         ("Priced in reality", "₹999–₹2,999/month against ₹190 of monthly compute at high usage.")]
for i, (t, d) in enumerate(india):
    x = 0.7 + (i % 2) * 6.2
    y = 1.7 + (i // 2) * 1.6
    rrect(s, x, y, 5.9, 1.4, WHITE, BORDER, 0.12)
    oval(s, x + 0.25, y + 0.33, 0.55, 0.55, SAND)
    txt(s, str(i + 1), x + 0.25, y + 0.44, 0.55, 0.3, size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txt(s, t, x + 0.98, y + 0.18, 4.7, 0.3, size=13.5, bold=True, color=NAVY)
    txt(s, d, x + 0.98, y + 0.52, 4.75, 0.8, size=10.5, color=MUTED, spacing=1.25)
notes(s, "India-first is a wedge, not decoration. Every one of these is in the codebase.")

# ══ 13 · WHAT THE TESTS FOUND ═══════════════════════════════════════════════
s = slide()
eyebrow(s, "WHAT THE TESTS FOUND")
title(s, "Two production bugs, caught before a judge did.", y=0.78, size=31, h=0.7)
rrect(s, 0.7, 1.8, 5.95, 2.1, WHITE, ROSE_LN, 0.12)
txt(s, "42P17 · RLS infinite recursion", 1.0, 2.0, 5.4, 0.3, size=14, bold=True, color=ROSE_TX)
txt(s, "Multi-tenancy policies queried the employees table from inside a policy on that same table. "
       "Postgres aborted every query, so every database write in production was silently failing "
       "while the interface reported success.",
    1.0, 2.4, 5.4, 1.4, size=11.5, color=INK, spacing=1.28)
rrect(s, 6.7, 1.8, 5.95, 2.1, WHITE, GOLD, 0.12)
txt(s, "22P02 · invalid numeric input", 7.0, 2.0, 5.4, 0.3, size=14, bold=True, color=AMB_TX)
txt(s, "A word like “negative” reaching a numeric sentiment column made Postgres reject the entire "
       "insert. The ticket never saved — and the interface still showed success.",
    7.0, 2.4, 5.4, 1.4, size=11.5, color=INK, spacing=1.28)
rrect(s, 0.7, 4.15, 11.95, 2.0, NAVY, None, 0.12)
txt(s, "Both were invisible from the UI — the code degrades gracefully instead of crashing.",
    1.05, 4.42, 11.2, 0.35, size=15, bold=True, color=WHITE)
txt(s, "Fixed with SECURITY DEFINER functions — the correct Postgres pattern for self-referential membership\n"
       "tables — and input coercion, each with a regression test so neither can return silently.",
    1.05, 4.88, 11.2, 0.8, size=12.5, color=ICE, spacing=1.35)
txt(s, "We would rather find these ourselves than have them found for us.",
    1.05, 5.68, 11.2, 0.35, size=12.5, italic=True, color=GOLD_LT)
notes(s, "Strongest engineering story: wrote tests, tests exposed a real outage, fixed it properly.")

# ══ 14 · HONEST VALIDATION ══════════════════════════════════════════════════
s = slide()
eyebrow(s, "WHAT'S PROVEN AND WHAT ISN'T")
title(s, "We publish the gaps too.", y=0.78, size=34, h=0.65)
vrows = [("Company Brain engine", "Unit-tested plus a reproducible scenario. Same input, same output.", MINT_BG, MINT_TX),
         ("The closed loop, in-app", "Visible in CareFlow and SalesFlow; context injected into the live prompt.", MINT_BG, MINT_TX),
         ("Live Groq API run", "Four real calls — latency, tokens and cost measured and logged.", MINT_BG, MINT_TX),
         ("LLM judgement calls", "Bias scores and sentiment are drafts for human review, not certified results.", AMB_BG, AMB_TX),
         ("Real customer data", "Inputs are still project-authored. Not yet run on a real business's records.", ROSE_BG, ROSE_TX)]
for i, (t, d, bg, fg) in enumerate(vrows):
    y = 1.75 + i * 0.92
    rrect(s, 0.7, y, 11.95, 0.8, bg, None, 0.1)
    oval(s, 1.0, y + 0.28, 0.24, 0.24, fg)
    txt(s, t, 1.45, y + 0.24, 3.6, 0.32, size=13.5, bold=True, color=fg)
    txt(s, d, 5.15, y + 0.26, 7.3, 0.32, size=12, color=INK)
txt(s, "The red row stays red until it's earned. That's why you can trust the green ones.",
    0.7, 6.45, 11.95, 0.4, size=13.5, bold=True, italic=True, color=NAVY, align=PP_ALIGN.CENTER)
notes(s, "Judges probe overclaiming. Pre-empting it is a strength.")

# ══ 15 · CLOSING ════════════════════════════════════════════════════════════
s = slide(NAVY); dark_decor(s)
txt(s, "Six agents that don't share memory\nare still six silos.", 0.9, 1.5, 11, 1.8,
    size=38, bold=True, color=WHITE, font=H, spacing=1.28)
txt(s, "NexFlow AI gives them one.", 0.9, 3.3, 11, 0.7, size=31, italic=True, color=GOLD_LT, font=H)
for i, (v, l) in enumerate([("7", "AI systems"), ("128", "tests"), ("₹0.19", "per run"), ("MIT", "open source")]):
    x = 0.9 + i * 3.02
    rrect(s, x, 4.35, 2.72, 1.15, PANEL, TEAL, 0.12)
    txt(s, v, x, 4.5, 2.72, 0.6, size=26, bold=True, color=GOLD_LT, font=H, align=PP_ALIGN.CENTER)
    txt(s, l, x, 5.08, 2.72, 0.32, size=11.5, color=ICE, align=PP_ALIGN.CENTER)
txt(s, "Mukul · IIIT-Delhi    |    with Supriya Bansal · IIT-Delhi", 0.9, 5.85, 11, 0.35,
    size=13.5, bold=True, color=WHITE)
txt(s, "nexflow-ai-india.netlify.app    ·    github.com/Mukul07777/hireflow-ai", 0.9, 6.25, 11, 0.35,
    size=12, color=GOLD_LT)
txt(s, "Thank you.", 9.9, 1.5, 2.7, 0.5, size=18, italic=True, color=STEEL, font=H, align=PP_ALIGN.RIGHT)
notes(s, "Land the one-liner, pause, then invite questions.")

out = os.path.join(ROOT, "NexFlow-AI-Presentation.pptx")
prs.save(out)
print("Created:", out, "|", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
